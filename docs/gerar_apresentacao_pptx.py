"""Gera o .pptx editável a partir do mesmo HTML da apresentação.

O PPTX NÃO é feito de imagens dos slides: cada texto vira caixa de texto, cada
tabela vira tabela do PowerPoint e cada gráfico entra como figura. Assim dá para
ajustar detalhes depois sem refazer nada.

Identidade GREI: verde #2F5E4B, creme #F4FFEB, do manual de marca.
"""
import base64
import io
import json
import pathlib
import re

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt
from PIL import Image

AQUI = pathlib.Path(__file__).resolve().parent / "apresentacao_assets"
REPO = pathlib.Path(__file__).resolve().parents[1]
HTML = REPO / "docs/apresentacao_mercado.html"
SAIDA = REPO / "docs/apresentacao_mercado.pptx"

VERDE = RGBColor(0x2F, 0x5E, 0x4B)
CREME = RGBColor(0xF4, 0xFF, 0xEB)
PAPEL = RGBColor(0xFB, 0xFE, 0xF6)
TINTA = RGBColor(0x16, 0x24, 0x1D)
TINTA2 = RGBColor(0x35, 0x47, 0x3E)
CINZA = RGBColor(0x64, 0x78, 0x6C)
LINHA = RGBColor(0xDD, 0xE7, 0xDA)
QUENTE = RGBColor(0xB6, 0x5A, 0x2B)
BOM = RGBColor(0x2F, 0x7D, 0x57)
CARTAO = RGBColor(0xFF, 0xFF, 0xFF)
ACENTO_FRACO = RGBColor(0xE4, 0xF0, 0xE5)

TITULO = "Exo 2"          # fonte da marca; cai para a alternativa se ausente
CORPO = "Poppins"
MONO = "Consolas"

L, T = Inches(0.85), Inches(0.62)
LARG = Inches(13.333) - 2 * L


def limpa(txt):
    return re.sub(r"\s+", " ", txt or "").strip()


def add_caixa(slide, x, y, cx, cy):
    tb = slide.shapes.add_textbox(x, y, cx, cy)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def escreve(par, no, tam, cor, negrito_cor=None, fonte=CORPO):
    """Transfere o conteúdo de um nó HTML preservando o negrito.

    O espaço de borda vem do HTML e nao e inventado: acrescentar um espaco depois
    de cada trecho produzia "reativo , que", com espaco antes da virgula.
    """
    partes = no.contents if hasattr(no, "contents") else [no]
    for parte in partes:
        nome = getattr(parte, "name", None)
        bruto = parte.get_text() if nome else str(parte)
        texto = re.sub(r"\s+", " ", bruto)
        if not texto.strip():
            continue
        r = par.add_run()
        r.text = texto
        r.font.size = Pt(tam)
        r.font.name = fonte
        forte = nome in ("strong", "b")
        r.font.bold = forte
        r.font.color.rgb = (negrito_cor or TINTA) if forte else cor


def bloco_texto(slide, x, y, cx, itens, tam=13, cor=TINTA2, espaco=6):
    """itens: lista de (no_html, marcador) — marcador None para parágrafo."""
    tf = add_caixa(slide, x, y, cx, Inches(0.4))
    primeiro = True
    for no, marcador in itens:
        p = tf.paragraphs[0] if primeiro else tf.add_paragraph()
        primeiro = False
        p.space_after = Pt(espaco)
        p.line_spacing = 1.25
        if marcador:
            r = p.add_run()
            r.text = marcador + "  "
            r.font.size = Pt(tam)
            r.font.name = MONO
            r.font.bold = True
            r.font.color.rgb = VERDE
        escreve(p, no, tam, cor)
    return tf


def desenha_fundo(slide, cor):
    fundo = slide.background.fill
    fundo.solid()
    fundo.fore_color.rgb = cor


def logo_bytes():
    """O PowerPoint nao recolore imagem como a mascara CSS do HTML faz, entao as
    duas variantes ja vem tingidas: creme para a capa verde, verde para o resto."""
    figs = json.loads((AQUI / "figs.json").read_text())
    return ((AQUI / "logo_creme.png").read_bytes(),
            (AQUI / "logo_verde.png").read_bytes(),
            figs["_logo_ratio"])


def img_bytes(src):
    return base64.b64decode(src.split(",", 1)[1])


def main():
    sopa = BeautifulSoup(HTML.read_text(), "lxml")
    slides = sopa.select("section.slide")
    logo_creme, logo_verde, logo_ratio = logo_bytes()

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    branco = prs.slide_layouts[6]

    for i, sec in enumerate(slides):
        s = prs.slides.add_slide(branco)
        capa = "cover" in (sec.get("class") or [])
        desenha_fundo(s, VERDE if capa else PAPEL)
        y = T

        # --- logotipo ---------------------------------------------------
        if capa:
            larg = Inches(2.6)
            s.shapes.add_picture(io.BytesIO(logo_creme), L, y,
                                 width=larg, height=Emu(int(larg / logo_ratio)))
            y += Emu(int(larg / logo_ratio)) + Inches(0.45)
        else:
            larg = Inches(0.8)
            s.shapes.add_picture(io.BytesIO(logo_verde),
                                 prs.slide_width - L - larg,
                                 prs.slide_height - Inches(0.78),
                                 width=larg, height=Emu(int(larg / logo_ratio)))

        # --- eyebrow ----------------------------------------------------
        eb = sec.select_one(".eyebrow")
        if eb:
            tf = add_caixa(s, L, y, LARG, Inches(0.26))
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = limpa(eb.get_text()).upper()
            r.font.size = Pt(10)
            r.font.name = MONO
            r.font.bold = True
            r.font.color.rgb = CREME if capa else VERDE
            y += Inches(0.36)

        # --- título -----------------------------------------------------
        h = sec.select_one("h1, h2")
        if h:
            grande = h.name == "h1"
            tf = add_caixa(s, L, y, LARG, Inches(0.9))
            p = tf.paragraphs[0]
            p.line_spacing = 1.05
            r = p.add_run()
            r.text = limpa(h.get_text())
            r.font.size = Pt(34 if grande else 27)
            r.font.name = TITULO
            r.font.bold = True
            r.font.color.rgb = CREME if capa else TINTA
            linhas = 1 + len(limpa(h.get_text())) // 52
            y += Inches(0.52) * linhas + Inches(0.22)

        corpo_y = y
        col_dir_x = L + LARG / 2 + Inches(0.25)
        meia = LARG / 2 - Inches(0.25)

        # --- extração em blocos, na ordem do documento -------------------
        # Percorrer `.inner` em ordem evita a perda de conteúdo que acontece ao
        # selecionar só filhos diretos: parágrafos dentro de colunas não são
        # filhos de `.inner`, e sumiam.
        inner = sec.select_one(".inner")
        CARTAO_CLS = {"card", "panel", "node"}

        def eh_cartao(el):
            return bool(CARTAO_CLS & set(el.get("class") or []))

        blocos = []

        def coleta(el):
            for filho in el.find_all(recursive=False):
                cls = set(filho.get("class") or [])
                if filho.name in ("h1", "h2") or "eyebrow" in cls:
                    continue
                if filho.name == "figure":
                    blocos.append(("figura", filho)); continue
                # So a propria <table> conta. Testar "contem tabela" fazia a
                # COLUNA inteira ser tratada como tabela, e os paragrafos ao lado
                # dela desapareciam.
                if filho.name == "table":
                    blocos.append(("tabela", filho)); continue
                if eh_cartao(filho):
                    blocos.append(("cartao", filho)); continue
                if "stat" in cls:
                    blocos.append(("stat", filho)); continue
                if "hero-number" in cls:
                    blocos.append(("heroi", filho)); continue
                if filho.name == "p":
                    blocos.append(("p", filho)); continue
                if filho.name == "ul":
                    blocos.append(("lista", filho)); continue
                if filho.name in ("div", "span"):
                    coleta(filho); continue

        coleta(inner)

        textos = [(n, k) for k, n in blocos if k in ("p", "lista")]
        cartoes = [n for k, n in blocos if k == "cartao"]
        stats = [n for k, n in blocos if k == "stat"]
        figura = next((n for k, n in blocos if k == "figura"), None)
        tabela = next((n for k, n in blocos if k == "tabela"), None)
        heroi = next((n for k, n in blocos if k == "heroi"), None)

        # --- figura ocupa o slide inteiro -------------------------------
        if figura:
            dados = img_bytes(figura.select_one("img")["src"])
            im = Image.open(io.BytesIO(dados))
            legenda = figura.select_one("figcaption")
            notas = [n for n, k in textos if k == "p"]
            alt_texto = Inches(0.34) * (1 + len(notas)) + Inches(0.5)
            disp_h = prs.slide_height - corpo_y - alt_texto - Inches(0.5)
            larg_fig = min(LARG, Emu(int(disp_h * im.width / im.height)))
            alt_fig = Emu(int(larg_fig * im.height / im.width))
            s.shapes.add_picture(io.BytesIO(dados),
                                 L + (LARG - larg_fig) // 2, corpo_y,
                                 width=larg_fig, height=alt_fig)
            tf = add_caixa(s, L, corpo_y + alt_fig + Inches(0.14), LARG, Inches(0.5))
            primeiro = True
            for no in ([legenda] if legenda else []) + notas:
                p = tf.paragraphs[0] if primeiro else tf.add_paragraph()
                primeiro = False
                p.space_after = Pt(3); p.line_spacing = 1.2
                escreve(p, no, 10, CINZA)
            continue

        def desenha_cartoes(y0):
            """Desenha os cartoes a partir de y0 e devolve o y logo abaixo."""
            if not cartoes:
                return y0
            largura_total = len(cartoes) > 3
            n_col = 3 if largura_total else 1
            larg_c = (LARG / 3 - Inches(0.18)) if largura_total else meia
            alt_c = Inches(1.95) if largura_total else Inches(1.62)
            for n, c in enumerate(cartoes):
                if n_col == 1:
                    cx, cy = col_dir_x, y0 + n * (alt_c + Inches(0.22))
                else:
                    cx = L + (n % 3) * (LARG / 3)
                    cy = y0 + (n // 3) * (alt_c + Inches(0.25))
                classes = set(c.get("class") or [])
                cor_fundo = (RGBColor(0xF8, 0xEA, 0xE0) if "before" in classes else
                             RGBColor(0xE0, 0xF1, 0xE6) if "after" in classes else
                             VERDE if "hub" in classes else CARTAO)
                escuro = "hub" in classes
                cxa = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy,
                                         larg_c, alt_c)
                cxa.adjustments[0] = 0.06
                cxa.fill.solid(); cxa.fill.fore_color.rgb = cor_fundo
                cxa.line.color.rgb = LINHA
                cxa.shadow.inherit = False
                tf = cxa.text_frame
                tf.word_wrap = True
                tf.margin_left = tf.margin_right = Inches(0.15)
                tf.margin_top = tf.margin_bottom = Inches(0.11)
                tf.vertical_anchor = MSO_ANCHOR.TOP
                primeiro = True
                for filho in c.select(".tag, h3, p, .nm, .rl, li"):
                    par = tf.paragraphs[0] if primeiro else tf.add_paragraph()
                    primeiro = False
                    par.space_after = Pt(3)
                    par.alignment = PP_ALIGN.LEFT
                    cls = set(filho.get("class") or [])
                    if "tag" in cls:
                        r = par.add_run(); r.text = limpa(filho.get_text()).upper()
                        r.font.size = Pt(8); r.font.name = MONO
                        r.font.color.rgb = CREME if escuro else CINZA
                    elif filho.name == "h3" or "nm" in cls:
                        r = par.add_run(); r.text = limpa(filho.get_text())
                        r.font.size = Pt(12.5); r.font.name = TITULO
                        r.font.bold = True
                        r.font.color.rgb = CREME if escuro else TINTA
                    else:
                        if filho.name == "li":
                            r = par.add_run(); r.text = "\u25cb  "
                            r.font.size = Pt(10); r.font.name = MONO
                            r.font.color.rgb = CREME if escuro else VERDE
                        escreve(par, filho, 10, CREME if escuro else TINTA2,
                                negrito_cor=CREME if escuro else TINTA)
            linhas_c = 1 if n_col == 1 else (len(cartoes) + 2) // 3
            usados = (len(cartoes) if n_col == 1 else linhas_c)
            return y0 + usados * (alt_c + Inches(0.24))

        # --- tabela ------------------------------------------------------
        if tabela is not None:
            # Os cartoes vem antes da tabela na leitura; desenha-los primeiro
            # evita que o ramo da tabela encerre o slide e os descarte.
            corpo_y = desenha_cartoes(corpo_y)
            cabec = [limpa(th.get_text()) for th in tabela.select("thead th")]
            corpo_linhas = tabela.select("tbody tr")
            dados_lin = [[limpa(td.get_text()) for td in tr.select("td")]
                         for tr in corpo_linhas]
            nl, nc = len(dados_lin) + 1, max(1, len(cabec))
            alt = min(Inches(4.9), Inches(0.38) * nl + Inches(0.15))
            forma = s.shapes.add_table(nl, nc, L, corpo_y, LARG, alt)
            t = forma.table
            t.first_row = False
            for j, txt in enumerate(cabec):
                cel = t.cell(0, j); cel.text = txt
                r0 = cel.text_frame.paragraphs[0].runs[0]
                r0.font.size = Pt(9.5); r0.font.name = MONO
                r0.font.bold = True; r0.font.color.rgb = CINZA
                cel.fill.solid(); cel.fill.fore_color.rgb = PAPEL
                cel.vertical_anchor = MSO_ANCHOR.MIDDLE
            for k, lin in enumerate(dados_lin, start=1):
                destaque = "hi" in (corpo_linhas[k - 1].get("class") or [])
                for j in range(nc):
                    cel = t.cell(k, j)
                    cel.text = lin[j] if j < len(lin) else ""
                    pr = cel.text_frame.paragraphs[0]
                    if pr.runs:
                        pr.runs[0].font.size = Pt(10.5)
                        pr.runs[0].font.name = CORPO
                        pr.runs[0].font.color.rgb = TINTA2
                    cel.fill.solid()
                    cel.fill.fore_color.rgb = ACENTO_FRACO if destaque else CARTAO
                    cel.vertical_anchor = MSO_ANCHOR.MIDDLE
            corpo_y += alt + Inches(0.3)
            notas = [n for n, k in textos if k == "p"]
            if notas:
                bloco_texto(s, L, corpo_y, LARG, [(n, None) for n in notas],
                            tam=10.5, cor=CINZA)
            continue

        # --- números em destaque ----------------------------------------
        if heroi is not None:
            tf = add_caixa(s, L, corpo_y, LARG, Inches(1.2))
            p = tf.paragraphs[0]
            for span in heroi.select(".big, .arrow"):
                txt = limpa(span.get_text())
                r = p.add_run(); r.text = txt + "   "
                r.font.size = Pt(52); r.font.name = TITULO; r.font.bold = True
                r.font.color.rgb = (QUENTE if txt == "337"
                                    else BOM if txt == "0" else CINZA)
            rot = heroi.select_one(".stat .l")
            if rot is not None:
                p2 = tf.add_paragraph(); p2.space_before = Pt(6)
                escreve(p2, rot, 11, CINZA)
            corpo_y += Inches(1.5)

        # --- texto e cartões --------------------------------------------
        # Com ate tres cartoes eles ficam na metade direita e o texto na
        # esquerda. Com mais, os cartoes ocupam a largura toda em tres colunas, e
        # ai o texto tem de vir DEPOIS deles, senao fica desenhado por baixo.
        largura_total = len(cartoes) > 3
        n_col = 3 if largura_total else 1
        larg_c = (LARG / 3 - Inches(0.18)) if largura_total else meia
        alt_c = Inches(1.95) if largura_total else Inches(1.62)
        y_cartoes = corpo_y
        if largura_total and cartoes:
            n_linhas = (len(cartoes) + 2) // 3
            corpo_y += n_linhas * (alt_c + Inches(0.25)) + Inches(0.15)

        largura_texto = meia if (cartoes and not largura_total) else LARG
        if textos:
            itens = []
            for no, tipo in textos:
                if tipo == "p":
                    itens.append((no, None))
                else:
                    passo = "steps" in (no.get("class") or [])
                    for k, li in enumerate(no.select("li"), start=1):
                        itens.append((li, f"{k}." if passo else "\u25cb"))
            bloco_texto(s, L, corpo_y, largura_texto, itens, tam=12,
                        cor=CREME if capa else TINTA2)

        if stats:
            passo_x = LARG / max(1, len(stats))
            colx = L
            base_y = prs.slide_height - Inches(2.2) if capa else corpo_y + Inches(0.1)
            for st in stats:
                tf = add_caixa(s, colx, base_y, passo_x - Inches(0.3), Inches(1.3))
                p = tf.paragraphs[0]
                r = p.add_run(); r.text = limpa(st.select_one(".n").get_text())
                r.font.size = Pt(28); r.font.name = TITULO; r.font.bold = True
                r.font.color.rgb = CREME if capa else VERDE
                p2 = tf.add_paragraph(); p2.space_before = Pt(3)
                escreve(p2, st.select_one(".l"), 9.5, CREME if capa else CINZA)
                colx += passo_x

        desenha_cartoes(y_cartoes)

        # numeração do slide
        tf = add_caixa(s, prs.slide_width - Inches(2.3),
                       prs.slide_height - Inches(0.45), Inches(0.6), Inches(0.25))
        r = tf.paragraphs[0].add_run(); r.text = f"{i + 1}"
        r.font.size = Pt(9); r.font.name = MONO
        r.font.color.rgb = CREME if capa else CINZA
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

    prs.save(SAIDA)
    print(f"{SAIDA} gravado, {len(slides)} slides, "
          f"{SAIDA.stat().st_size / 1e6:.2f} MB")


if __name__ == "__main__":
    main()
